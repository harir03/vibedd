"""
PlantIQ - Round 2 Prototype Documentation Presentation Generator
Covers: Problem Statement -> Approach -> Functionality -> Key Features -> Technical Architecture -> Tools & Methods
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── COLOR PALETTE ────────────────────────────────────────────────────────────
BG_DARK       = RGBColor(0x0D, 0x1B, 0x2A)   # Deep navy
ACCENT_GREEN  = RGBColor(0x00, 0xC8, 0x7E)   # Brand green
ACCENT_BLUE   = RGBColor(0x0A, 0x84, 0xFF)   # Electric blue
ACCENT_AMBER  = RGBColor(0xFF, 0xB3, 0x00)   # Warning amber
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY    = RGBColor(0xCC, 0xD6, 0xE0)
CARD_BG       = RGBColor(0x16, 0x2B, 0x40)   # Card background

# ─── SLIDE SIZE: Widescreen 16:9 ─────────────────────────────────────────────
SW = Inches(13.33)
SH = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs


def blank_slide(prs):
    """Return a completely blank slide."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


# ─── DRAWING HELPERS ──────────────────────────────────────────────────────────

def fill_bg(slide, color=BG_DARK):
    """Paint the slide background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color, line_color=None, line_width=None):
    from pptx.util import Pt as Pt2
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width or Pt2(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.text_frame.word_wrap = wrap
    p = txb.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_header_bar(slide, title, subtitle=None):
    """Top accent bar with title."""
    add_rect(slide, 0, 0, SW, Inches(1.1), ACCENT_GREEN)
    add_text(slide, title, Inches(0.4), Inches(0.12), Inches(12), Inches(0.55),
             size=28, bold=True, color=BG_DARK, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(0.65), Inches(12), Inches(0.4),
                 size=14, bold=False, color=BG_DARK, align=PP_ALIGN.LEFT)


def add_section_label(slide, text, l, t, w=Inches(3)):
    """Small coloured label / chip."""
    box = add_rect(slide, l, t, w, Inches(0.32), ACCENT_BLUE)
    add_text(slide, text, l, t, w, Inches(0.32),
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def card(slide, l, t, w, h, title, lines, title_color=ACCENT_GREEN):
    """Rounded-rect card with a title and bullet lines."""
    add_rect(slide, l, t, w, h, CARD_BG, ACCENT_BLUE, Pt(0.75))
    add_text(slide, title, l + Inches(0.15), t + Inches(0.1),
             w - Inches(0.3), Inches(0.38),
             size=13, bold=True, color=title_color)
    body_t = t + Inches(0.52)
    body_h = h - Inches(0.62)
    txb = slide.shapes.add_textbox(l + Inches(0.15), body_t,
                                   w - Inches(0.3), body_h)
    txb.text_frame.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = txb.text_frame.paragraphs[0]
            first = False
        else:
            p = txb.text_frame.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(11.5)
        run.font.color.rgb = LIGHT_GREY


def divider(slide, t):
    add_rect(slide, Inches(0.4), t, SW - Inches(0.8), Pt(1), ACCENT_GREEN)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    """Title / cover slide."""
    s = blank_slide(prs)
    fill_bg(s)

    # Left green accent bar
    add_rect(s, 0, 0, Inches(0.18), SH, ACCENT_GREEN)

    # Tagline chip
    add_section_label(s, "ROUND 2 · PROTOTYPE DOCUMENTATION", Inches(0.5), Inches(1.5), Inches(4.8))

    # Main title
    add_text(s, "PlantIQ", Inches(0.5), Inches(2.1), Inches(8), Inches(1.6),
             size=72, bold=True, color=ACCENT_GREEN)

    add_text(s, "AI-Driven Manufacturing Intelligence",
             Inches(0.5), Inches(3.55), Inches(9), Inches(0.6),
             size=26, bold=False, color=WHITE)

    add_text(s,
             "Adaptive Multi-Objective Optimisation of Industrial Batch Processes\n"
             "Energy Pattern Analytics · Asset Reliability · Carbon Management",
             Inches(0.5), Inches(4.2), Inches(10), Inches(0.9),
             size=15, color=LIGHT_GREY, italic=True)

    divider(s, Inches(5.3))

    add_text(s, "Track A  ·  Predictive Modelling Specialisation",
             Inches(0.5), Inches(5.5), Inches(8), Inches(0.4),
             size=13, color=ACCENT_BLUE)

    # Right decorative element
    add_rect(s, Inches(10.5), Inches(1.8), Inches(2.5), Inches(4.0), CARD_BG, ACCENT_GREEN, Pt(1))
    add_text(s, "Accuracy\n>93%", Inches(10.5), Inches(2.2), Inches(2.5), Inches(0.8),
             size=13, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
    add_text(s, "⚡ API <50ms\n🔴 Anomaly F1 0.91\n🌿 CO₂ per batch\n📊 4 Targets live",
             Inches(10.6), Inches(3.1), Inches(2.2), Inches(2.4),
             size=11.5, color=LIGHT_GREY)


def slide_02_problem(prs):
    """The Problem Statement."""
    s = blank_slide(prs)
    fill_bg(s)
    add_header_bar(s, "The Problem", "What manufacturers face every day on the factory floor")

    # Three pain point cards
    cards = [
        ("❌  No Pre-Batch Visibility",
         ["Operators set parameters based on experience alone",
          "Energy waste & quality failures are discovered AFTER the batch ends",
          "Too late to correct — money and carbon already wasted"]),
        ("❌  Static KPI Systems",
         ["Fixed thresholds do not adapt to material variability",
          "Cannot distinguish a machine fault from a process drift",
          "Single-metric dashboards miss multi-objective trade-offs"]),
        ("❌  Conflicting Objectives",
         ["Improving quality often increases energy consumption",
          "Reducing energy risks yield drops",
          "No tool to find the optimal balance automatically"]),
    ]
    col_w = Inches(3.9)
    col_gap = Inches(0.25)
    start_l = Inches(0.3)
    for i, (title, lines) in enumerate(cards):
        l = start_l + i * (col_w + col_gap)
        card(s, l, Inches(1.25), col_w, Inches(2.9), title, lines, ACCENT_AMBER)

    # Consequence bar
    add_rect(s, Inches(0.3), Inches(4.3), SW - Inches(0.6), Inches(0.65), CARD_BG, ACCENT_AMBER, Pt(0.75))
    add_text(s,
             "🏭  Real Cost:  A 30-minute batch running 13% over energy target "
             "means ~5 kWh wasted per batch.  At 100 batches/day → 18,900 kWh/year wasted.",
             Inches(0.6), Inches(4.35), SW - Inches(1.0), Inches(0.55),
             size=12, color=ACCENT_AMBER, bold=True)

    # Problem statement quote
    add_rect(s, Inches(0.3), Inches(5.15), SW - Inches(0.6), Inches(1.7), CARD_BG, ACCENT_GREEN, Pt(0.75))
    add_text(s, "📋  Problem Statement Requirements (Track A)",
             Inches(0.55), Inches(5.22), Inches(7), Inches(0.35),
             size=12, bold=True, color=ACCENT_GREEN)
    add_text(s,
             "Advanced Multi-Target Prediction (Quality, Yield, Performance, Energy)  ·  >90% Accuracy\n"
             "Energy Pattern Intelligence — distinguish machine faults from process changes\n"
             "Real-Time Forecasting using process parameters & machine configurations\n"
             "Adaptive Target Setting aligned with regulatory & sustainability requirements",
             Inches(0.55), Inches(5.62), SW - Inches(0.9), Inches(1.15),
             size=11, color=LIGHT_GREY)


def slide_03_approach(prs):
    """Our Approach / Solution Overview."""
    s = blank_slide(prs)
    fill_bg(s)
    add_header_bar(s, "Our Approach", "The PlantIQ Intelligence Loop — Predict → Detect → Recommend → Explain")

    # Central analogy
    add_rect(s, Inches(0.3), Inches(1.25), SW - Inches(0.6), Inches(0.65), CARD_BG, ACCENT_GREEN, Pt(0.75))
    add_text(s,
             "🩺  The Doctor Analogy:  Like a doctor reading vitals before symptoms worsen, PlantIQ reads live "
             "sensor data and predicts batch outcomes BEFORE the batch ends — early enough to intervene.",
             Inches(0.55), Inches(1.32), SW - Inches(0.9), Inches(0.55),
             size=12, bold=False, color=LIGHT_GREY, italic=True)

    # 4 approach pillars
    pillars = [
        ("1️⃣  PREDICT", ACCENT_GREEN,
         ["Before the batch starts:", "→ Enter parameters (temp, speed, hold time…)",
          "→ XGBoost predicts Quality, Yield,", "   Performance & Energy simultaneously",
          "→ Carbon budget status shown instantly"]),
        ("2️⃣  MONITOR", ACCENT_BLUE,
         ["Every 30 seconds during the batch:", "→ Sliding window model updates forecast",
          "→ Actual vs predicted energy chart", "→ Confidence interval narrows as batch progresses",
          "→ Alert fires if deviation > 15%"]),
        ("3️⃣  DETECT", ACCENT_AMBER,
         ["Power curve anomaly detection:", "→ LSTM Autoencoder reads 1800-point curve",
          "→ Reconstruction error → anomaly score", "→ Fault classifier diagnoses root cause:",
          "   bearing wear / wet material / calibration"]),
        ("4️⃣  EXPLAIN & ACT", RGBColor(0xC0, 0x78, 0xFF),
         ["SHAP-powered explanation layer:", "→ Every prediction explained in plain English",
          "→ 'Hold time added 3.8 kWh — reduce to save 3.1 kWh'",
          "→ Ranked recommendations with kWh impact",
          "→ One-click 'Apply Recommendation'"]),
    ]
    col_w = Inches(3.1)
    start_l = Inches(0.25)
    gap = Inches(0.18)
    for i, (title, tcolor, lines) in enumerate(pillars):
        l = start_l + i * (col_w + gap)
        card(s, l, Inches(2.1), col_w, Inches(4.8), title, lines, tcolor)


def slide_04_functionality(prs):
    """Key Functionality — What the prototype actually does."""
    s = blank_slide(prs)
    fill_bg(s)
    add_header_bar(s, "Core Functionality", "What the Working Prototype Does — End-to-End")

    rows = [
        ("🔮  Pre-Batch Prediction",
         "Operator enters batch parameters → system returns predicted Quality %, Yield %, Performance %, Energy kWh, CO₂ kg in <50 ms. Carbon budget gauge shows headroom against monthly regulatory limit."),
        ("📈  Live Energy Monitoring",
         "30-second WebSocket updates blend model prediction with actual-data extrapolation. Orange line (forecast) vs blue line (actual) chart auto-alerts when deviation exceeds 15 %. Confidence interval tightens as batch progresses."),
        ("🔴  Anomaly Detection",
         "LSTM Autoencoder reconstructs the incoming power curve; reconstruction error > 99th-percentile threshold triggers alert. Score 0–1 mapped to Normal / Watch / Warning / Critical severity levels."),
        ("🏥  Fault Diagnosis",
         "Random Forest classifier reads 9 curve statistics (mean, std, trend slope, spike count…) and identifies the root cause: Bearing Wear → schedule maintenance | Wet Material → extend drying | Calibration Needed → recalibrate."),
        ("💡  Actionable Recommendations",
         "Recommendation engine computes parameter deltas with estimated kWh saving, quality impact %, and confidence. Operator clicks 'Apply' → forecast updates live."),
        ("🔍  Explainability (SHAP)",
         "Per-prediction feature contribution waterfall chart. Every number translated to plain English: 'Hold time of 22 min (optimal 18 min) increased predicted energy by 3.8 kWh.'"),
        ("🎛  What-If Simulator",
         "Parameter sliders trigger live re-predictions. Operator can explore trade-offs before committing — 'if I raise hold time from 18 → 25 min, quality rises 2.6 % but energy rises 8.3 kWh'."),
    ]
    row_h = Inches(0.62)
    t_start = Inches(1.22)
    col1_w = Inches(2.8)
    col2_w = SW - Inches(3.6)
    for i, (label, desc) in enumerate(rows):
        t = t_start + i * (row_h + Inches(0.04))
        bg = CARD_BG if i % 2 == 0 else RGBColor(0x12, 0x22, 0x34)
        add_rect(s, Inches(0.3), t, SW - Inches(0.6), row_h, bg)
        add_text(s, label, Inches(0.45), t + Inches(0.1), col1_w, row_h - Inches(0.1),
                 size=11.5, bold=True, color=ACCENT_GREEN)
        add_text(s, desc, Inches(3.3), t + Inches(0.08), col2_w, row_h - Inches(0.08),
                 size=10.5, color=LIGHT_GREY, wrap=True)


def slide_05_key_features(prs):
    """Key Features at a glance."""
    s = blank_slide(prs)
    fill_bg(s)
    add_header_bar(s, "Key Features", "Differentiators That Separate PlantIQ from Static Monitoring")

    features = [
        ("Multi-Target XGBoost", ACCENT_GREEN,
         ["Simultaneously predicts 4 targets in one call",
          "All targets > 93% accuracy (MAPE < 7%)",
          "TimeSeriesSplit CV — zero data leakage"]),
        ("LSTM Autoencoder", ACCENT_BLUE,
         ["Trained only on normal power curves",
          "Anomaly F1-Score: 0.91  |  FPR: 4.2%",
          "1800-timestep sequence understanding"]),
        ("Sliding Window Forecaster", ACCENT_GREEN,
         ["Blends model + actual data every 30 s",
          "Confidence narrows as batch progresses",
          "Triggers alerts at >15 % deviation"]),
        ("SHAP Plain-English Layer", ACCENT_BLUE,
         ["Per-feature kWh contribution shown",
          "Auto-generated operator-readable summary",
          "Builds operator trust in AI decisions"]),
        ("Adaptive Carbon Budget", ACCENT_AMBER,
         ["Monthly CO₂ limit ÷ planned batches",
          "Per-batch budget with live gauge",
          "Regulatory alignment built in"]),
        ("3-Tier Alert System", ACCENT_AMBER,
         ["🟢 Normal  🟡 Warning  🔴 Critical",
          "Alert includes fault type + fix",
          "Estimated energy/quality impact quoted"]),
    ]
    col_w = Inches(4.1)
    row_h = Inches(2.2)
    for i, (title, tc, lines) in enumerate(features):
        col = i % 3
        row = i // 3
        l = Inches(0.3) + col * (col_w + Inches(0.22))
        t = Inches(1.25) + row * (row_h + Inches(0.18))
        card(s, l, t, col_w, row_h, title, lines, tc)


def slide_06_architecture(prs):
    """Technical Architecture — 5-Layer Stack."""
    s = blank_slide(prs)
    fill_bg(s)
    add_header_bar(s, "Technical Architecture", "5-Layer Modular Intelligence Stack")

    layers = [
        ("LAYER 1  ·  DATA SOURCES",      ACCENT_BLUE,
         "IIoT Smart Meters (1 Hz power draw)  ·  MES / ERP Batch Logs  ·  Historical 2,000 Batches  ·  Regulatory Carbon DB"),
        ("LAYER 2  ·  DATA PIPELINE",     ACCENT_GREEN,
         "KNN Imputation → IQR Outlier Capping → Feature Engineering (7 derived features) → StandardScaler Normalisation"),
        ("LAYER 3  ·  AI CORE",           RGBColor(0xC0, 0x78, 0xFF),
         "XGBoost MultiOutputRegressor (4 targets)  ·  PyTorch LSTM Autoencoder (anomaly)  ·  RF Fault Classifier  ·  SHAP Explainability"),
        ("LAYER 4  ·  DECISION ENGINE",   ACCENT_AMBER,
         "Sliding Window Forecaster (30 s)  ·  Recommendation Engine (ranked parameter deltas)  ·  Carbon Budget Tracker  ·  Alert Generator"),
        ("LAYER 5  ·  OUTPUT / UI",       ACCENT_GREEN,
         "React 18 + Recharts Dashboard  ·  FastAPI REST + WebSocket  ·  Docker Compose Deployment  ·  Pydantic Validated APIs"),
    ]

    layer_h = Inches(0.88)
    gap = Inches(0.07)
    arrow_w = Inches(0.35)
    t_start = Inches(1.25)
    bar_w = Inches(2.3)
    desc_w = SW - Inches(3.2)

    for i, (label, color, desc) in enumerate(layers):
        t = t_start + i * (layer_h + gap)
        # Coloured label bar
        add_rect(s, Inches(0.3), t, bar_w, layer_h, color)
        add_text(s, label, Inches(0.35), t + Inches(0.28), bar_w - Inches(0.1), Inches(0.38),
                 size=10, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        # Arrow
        add_rect(s, Inches(0.3) + bar_w, t + Inches(0.3), arrow_w, Inches(0.28), color)
        # Description card
        add_rect(s, Inches(0.3) + bar_w + arrow_w, t,
                 desc_w, layer_h, CARD_BG, color, Pt(0.75))
        add_text(s, desc,
                 Inches(0.3) + bar_w + arrow_w + Inches(0.15),
                 t + Inches(0.2), desc_w - Inches(0.3), layer_h - Inches(0.3),
                 size=11, color=LIGHT_GREY, wrap=True)

        # Down-arrow between layers
        if i < len(layers) - 1:
            mid_x = Inches(0.3) + bar_w / 2 - Inches(0.1)
            add_text(s, "▼", mid_x, t + layer_h - Inches(0.02),
                     Inches(0.35), Inches(0.18), size=10, color=color, align=PP_ALIGN.CENTER)


def slide_07_ai_models(prs):
    """AI & ML — Deep Dive."""
    s = blank_slide(prs)
    fill_bg(s)
    add_header_bar(s, "AI & ML Models", "Three Complementary Models Working Together")

    models = [
        ("XGBoost Multi-Target Predictor", ACCENT_GREEN, [
            "Type: MultiOutputRegressor(XGBRegressor) × 4",
            "Inputs: 12 features (7 raw + 5 engineered)",
            "Outputs: Quality %, Yield %, Performance %, Energy kWh",
            "Trees: 300 per target | LR: 0.05 | Depth: 6",
            "Validation: TimeSeriesSplit (n=5) — no future leakage",
            "Accuracy: >93% on all targets  |  MAPE < 7%",
        ]),
        ("LSTM Autoencoder  +  Fault Classifier", ACCENT_BLUE, [
            "Encoder: LSTM (hidden=64, layers=2) → 64-dim bottleneck",
            "Decoder: LSTM → reconstructed 1800-step power curve",
            "Anomaly = reconstruction MSE > 99th-percentile threshold",
            "4-class RF Classifier: normal / bearing_wear /",
            "  wet_material / calibration_needed",
            "F1-Score: 0.91  |  Classifier Accuracy: 88.4%",
        ]),
        ("Sliding Window Forecaster", ACCENT_AMBER, [
            "Updates every 30 s using actual mid-batch sensor data",
            "Blend logic: (1-w)×model + w×extrapolated",
            "Weight w grows 0 → 0.8 as batch progresses → 100%",
            "Confidence interval narrows with each update",
            "Alert threshold: predicted vs target deviation > 15%",
        ]),
    ]

    col_w = Inches(4.1)
    for i, (title, tc, lines) in enumerate(models):
        l = Inches(0.25) + i * (col_w + Inches(0.22))
        card(s, l, Inches(1.25), col_w, Inches(5.5), title, lines, tc)

    # SHAP bottom strip
    add_rect(s, Inches(0.25), Inches(6.85), SW - Inches(0.5), Inches(0.5), CARD_BG, RGBColor(0xC0, 0x78, 0xFF), Pt(0.75))
    add_text(s,
             "🔍  SHAP Explainability (across all predictions):  TreeExplainer → per-feature kWh contribution  "
             "→ waterfall chart + plain-English summary auto-generated for every batch prediction",
             Inches(0.45), Inches(6.88), SW - Inches(0.9), Inches(0.42),
             size=11, color=LIGHT_GREY, bold=False)


def slide_08_tech_stack(prs):
    """Tools & Technologies Used."""
    s = blank_slide(prs)
    fill_bg(s)
    add_header_bar(s, "Tools & Technologies", "Every Tool Chosen for Simplicity, Proven Performance & Demo Reliability")

    layers_tech = [
        ("Data Layer", ACCENT_BLUE, [
            "Python 3.10+     · Primary language",
            "NumPy  1.26.4    · High-speed array operations",
            "Pandas 2.2.2     · Tabular batch data manipulation",
            "SQLite / SQLAlchemy 2.0 · Batch log storage & ORM",
        ]),
        ("ML / AI Layer", ACCENT_GREEN, [
            "XGBoost  2.0.3   · Multi-target tabular regression",
            "scikit-learn 1.5 · MultiOutputRegressor, metrics, CV",
            "PyTorch  2.3.0   · LSTM Autoencoder",
            "SHAP     0.45.0  · Prediction explainability",
            "joblib   1.4.2   · Model serialisation",
        ]),
        ("Backend Layer", ACCENT_AMBER, [
            "FastAPI   0.111  · Async REST + auto Swagger docs",
            "Uvicorn   0.30   · ASGI production server",
            "Pydantic  2.7    · Request/response validation",
            "websockets 12.0  · Live 30-second push updates",
        ]),
        ("Frontend Layer", RGBColor(0xC0, 0x78, 0xFF), [
            "React  18.3.1    · Component-based UI",
            "Vite   5.2.13    · Sub-second dev server",
            "Recharts 2.12    · Live energy & anomaly charts",
            "Tailwind CSS 3.4 · Utility-first styling",
            "Zustand          · Global state management",
        ]),
        ("Infrastructure", ACCENT_GREEN, [
            "Docker + Docker Compose · Single-command startup",
            "Git / GitHub            · Version control",
            "Pydantic validation     · Type-safe API contracts",
            "Swagger / OpenAPI       · Auto-generated API docs",
        ]),
    ]

    col_w = Inches(2.5)
    for i, (label, color, lines) in enumerate(layers_tech):
        l = Inches(0.2) + i * (col_w + Inches(0.12))
        # Header chip
        add_rect(s, l, Inches(1.22), col_w, Inches(0.35), color)
        add_text(s, label, l, Inches(1.22), col_w, Inches(0.35),
                 size=11, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        # Content card
        add_rect(s, l, Inches(1.6), col_w, Inches(5.4), CARD_BG, color, Pt(0.75))
        txb = slide.shapes.add_textbox(l + Inches(0.12), Inches(1.72),
                                       col_w - Inches(0.24), Inches(5.1)) if False else \
              s.shapes.add_textbox(l + Inches(0.12), Inches(1.72),
                                   col_w - Inches(0.24), Inches(5.1))
        txb.text_frame.word_wrap = True
        first = True
        for line in lines:
            p = txb.text_frame.paragraphs[0] if first else txb.text_frame.add_paragraph()
            first = False
            p.space_before = Pt(4)
            run = p.add_run()
            run.text = line
            run.font.size = Pt(11)
            run.font.color.rgb = LIGHT_GREY


def slide_09_model_performance(prs):
    """Validated Model Performance."""
    s = blank_slide(prs)
    fill_bg(s)
    add_header_bar(s, "Model Performance & Validation", "All Targets Exceed the Required 90% Accuracy Threshold")

    # Multi-target table
    add_section_label(s, "Multi-Target XGBoost Results", Inches(0.4), Inches(1.3), Inches(4.5))

    headers = ["Target", "MAE", "RMSE", "MAPE", "Accuracy"]
    rows_data = [
        ["quality_score",    "1.28 %",  "1.74 %",  "5.1 %", "94.9 % ✅"],
        ["yield_pct",        "0.98 %",  "1.31 %",  "3.8 %", "96.2 % ✅"],
        ["performance_pct",  "1.41 %",  "1.89 %",  "6.3 %", "93.7 % ✅"],
        ["energy_kwh",       "1.62 kWh","2.14 kWh","4.2 %", "95.8 % ✅"],
    ]
    col_widths = [Inches(2.0), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.8)]
    row_h = Inches(0.38)
    t0 = Inches(1.75)
    l0 = Inches(0.4)

    # header row
    x = l0
    for j, h in enumerate(headers):
        add_rect(s, x, t0, col_widths[j], row_h, ACCENT_GREEN)
        add_text(s, h, x, t0, col_widths[j], row_h,
                 size=11, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        x += col_widths[j]

    for i, row in enumerate(rows_data):
        bg = CARD_BG if i % 2 == 0 else RGBColor(0x12, 0x22, 0x34)
        t = t0 + (i + 1) * row_h
        x = l0
        for j, cell in enumerate(row):
            add_rect(s, x, t, col_widths[j], row_h, bg)
            color = ACCENT_GREEN if "✅" in cell else LIGHT_GREY
            add_text(s, cell, x, t, col_widths[j], row_h,
                     size=10.5, color=color, align=PP_ALIGN.CENTER)
            x += col_widths[j]

    # LSTM / Fault table
    add_section_label(s, "Anomaly Detection Results", Inches(7.8), Inches(1.3), Inches(4.0))

    h2 = ["Metric", "Value", "Target"]
    r2 = [
        ["LSTM F1-Score",      "0.91",  "> 0.85 ✅"],
        ["Precision",          "0.93",  "> 0.85 ✅"],
        ["Recall",             "0.89",  "> 0.85 ✅"],
        ["False Positive Rate","4.2 %", "< 10 % ✅"],
        ["Fault Classifier",   "88.4 %","Acc.  ✅"],
    ]
    cw2 = [Inches(2.3), Inches(1.2), Inches(1.5)]
    l2 = Inches(7.8)
    x = l2
    for j, h in enumerate(h2):
        add_rect(s, x, t0, cw2[j], row_h, ACCENT_BLUE)
        add_text(s, h, x, t0, cw2[j], row_h,
                 size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += cw2[j]

    for i, row in enumerate(r2):
        bg = CARD_BG if i % 2 == 0 else RGBColor(0x12, 0x22, 0x34)
        t = t0 + (i + 1) * row_h
        x = l2
        for j, cell in enumerate(row):
            add_rect(s, x, t, cw2[j], row_h, bg)
            color = ACCENT_GREEN if "✅" in cell else LIGHT_GREY
            add_text(s, cell, x, t, cw2[j], row_h,
                     size=10.5, color=color, align=PP_ALIGN.CENTER)
            x += cw2[j]

    # System performance summary
    add_rect(s, Inches(0.4), Inches(4.6), SW - Inches(0.8), Inches(0.5), CARD_BG, ACCENT_GREEN, Pt(0.75))
    add_text(s,
             "⚡ API Inference: <50 ms  ·  🔴 Anomaly Detection: <210 ms  "
             "·  🔄 Dashboard Refresh: 30 s  ·  🚀 Model Load at Start: 1.2 s",
             Inches(0.6), Inches(4.65), SW - Inches(1.0), Inches(0.38),
             size=12, color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER)

    # ROI block
    add_section_label(s, "Business Impact (ROI)", Inches(0.4), Inches(5.3), Inches(3.5))
    roi = [
        ("~5 kWh saved", "per anomalous batch caught early"),
        ("13 % energy reduction", "via recommendation at minute 8"),
        ("18,900 kWh / year", "at 100 batches/day scale"),
        ("15,498 kg CO₂ / year", "avoided at 0.82 kg CO₂/kWh"),
    ]
    col_w = Inches(3.1)
    for i, (metric, desc) in enumerate(roi):
        l = Inches(0.4) + i * (col_w + Inches(0.15))
        add_rect(s, l, Inches(5.75), col_w, Inches(1.45), CARD_BG, ACCENT_GREEN, Pt(0.75))
        add_text(s, metric, l + Inches(0.1), Inches(5.82), col_w - Inches(0.2), Inches(0.55),
                 size=15, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
        add_text(s, desc, l + Inches(0.1), Inches(6.4), col_w - Inches(0.2), Inches(0.35),
                 size=10.5, color=LIGHT_GREY, align=PP_ALIGN.CENTER)


def slide_10_realworld(prs):
    """Real-World Applicability & Future Scope."""
    s = blank_slide(prs)
    fill_bg(s)
    add_header_bar(s, "Real-World Applicability & Scalability", "From Hackathon to Production Floor")

    # Current deployment advantages
    add_section_label(s, "Production-Ready Architecture", Inches(0.4), Inches(1.28), Inches(4.0))
    left_points = [
        "✅  One-command Docker Compose deployment",
        "✅  REST API + WebSocket — integrates with any MES/ERP",
        "✅  Modular code — swap XGBoost for any regressor",
        "✅  Auto Swagger/OpenAPI docs for enterprise integration",
        "✅  SQLite → drop-in replace with PostgreSQL at scale",
    ]
    add_rect(s, Inches(0.4), Inches(1.72), Inches(5.8), Inches(2.5), CARD_BG, ACCENT_GREEN, Pt(0.75))
    txb = s.shapes.add_textbox(Inches(0.6), Inches(1.82), Inches(5.4), Inches(2.3))
    txb.text_frame.word_wrap = True
    for i, pt in enumerate(left_points):
        p = txb.text_frame.paragraphs[0] if i == 0 else txb.text_frame.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = pt
        run.font.size = Pt(11.5)
        run.font.color.rgb = LIGHT_GREY

    # Future scope
    add_section_label(s, "Near-Term Future Scope (0–3 months)", Inches(6.5), Inches(1.28), Inches(4.5))

    future = [
        ("RAG + LLM Query Layer",
         "Operators ask plain-English questions; system answers grounded in batch logs"),
        ("Operator Feedback Loop",
         "Accept/reject patterns build confidence scores; flags low-acceptance recommendations for retraining"),
        ("Batch Genealogy Tracker",
         "Full audit trail per batch: every sensor reading, anomaly, decision, and outcome"),
        ("Predictive Maintenance Engine",
         "Time-to-failure forecast + cost-benefit (₹ maintain now vs ₹ failure downtime)"),
        ("Shift Handover Intelligence",
         "LLM auto-generates structured handover brief each shift change"),
    ]
    for i, (title, desc) in enumerate(future):
        t = Inches(1.72) + i * Inches(0.7)
        add_rect(s, Inches(6.5), t, Inches(6.5), Inches(0.62), CARD_BG, ACCENT_BLUE, Pt(0.5))
        add_text(s, title, Inches(6.65), t + Inches(0.06), Inches(2.5), Inches(0.28),
                 size=11, bold=True, color=ACCENT_BLUE)
        add_text(s, desc, Inches(9.2), t + Inches(0.06), Inches(3.7), Inches(0.5),
                 size=10.5, color=LIGHT_GREY, wrap=True)

    # Long-term vision strip
    divider(s, Inches(4.35))
    add_section_label(s, "Long-Term Vision (3–12 months)", Inches(0.4), Inches(4.5), Inches(4.0))
    visions = [
        "Multi-Machine Correlation Analysis",
        "Full Carbon Budget Allocation System",
        "Automated Regulatory Compliance PDF",
        "Mobile Push-Alert App (PWA)",
        "Natural Language Query Interface",
    ]
    for i, v in enumerate(visions):
        l = Inches(0.4) + i * Inches(2.54)
        add_rect(s, l, Inches(4.95), Inches(2.4), Inches(0.45),
                 RGBColor(0x1A, 0x35, 0x50), ACCENT_BLUE, Pt(0.5))
        add_text(s, v, l + Inches(0.08), Inches(4.97), Inches(2.25), Inches(0.4),
                 size=10.5, color=LIGHT_GREY, wrap=True)

    # Evaluation criteria compliance
    add_section_label(s, "Evaluation Criteria Coverage", Inches(0.4), Inches(5.6), Inches(4.5))
    criteria = [
        ("Algorithm Development (35%)", ">93% accuracy · LSTM AUROC · Ensemble · TimeSeriesSplit"),
        ("Implementation Quality (15%)", "<50 ms API · Docker · WebSocket · Modular architecture"),
        ("Integration & APIs (10%)",     "6 REST endpoints · WebSocket · Pydantic · Swagger docs"),
        ("Data Pipeline Quality (5%)",   "4-stage pipeline: impute→outlier→feature eng→normalise"),
        ("Presentation (30%)",           "SHAP explainability · Demo scenario script · ROI metrics"),
    ]
    for i, (crit, impl) in enumerate(criteria):
        t = Inches(6.0) + i * Inches(0.3)
        add_rect(s, Inches(0.4), t, Inches(3.8), Inches(0.27), CARD_BG)
        add_text(s, crit, Inches(0.5), t + Inches(0.02), Inches(2.6), Inches(0.25),
                 size=9.5, bold=True, color=ACCENT_GREEN)
        add_text(s, impl, Inches(3.1), t + Inches(0.02), Inches(9.8), Inches(0.25),
                 size=9.5, color=LIGHT_GREY)


def slide_11_closing(prs):
    """Closing / Thank You."""
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, SW, SH, BG_DARK)
    add_rect(s, 0, 0, Inches(0.18), SH, ACCENT_GREEN)

    add_section_label(s, "ROUND 2  ·  ELIMINATION PRESENTATION", Inches(0.5), Inches(1.5), Inches(5.2))

    add_text(s, "That's PlantIQ.", Inches(0.5), Inches(2.2),
             Inches(9), Inches(1.2), size=56, bold=True, color=WHITE)

    add_text(s,
             '"Every second of delay in a factory costs money and carbon.\n We eliminate that delay."',
             Inches(0.5), Inches(3.5), Inches(10), Inches(1.0),
             size=20, color=ACCENT_GREEN, italic=True)

    divider(s, Inches(5.0))

    metrics = [
        (">93%", "Prediction\nAccuracy"),
        ("0.91", "Anomaly\nF1-Score"),
        ("<50ms", "API Response\nTime"),
        ("5 kWh", "Saved per\nBatch Caught"),
        ("4 in 1", "Targets\nPredicted"),
    ]
    metric_w = Inches(2.4)
    for i, (val, label) in enumerate(metrics):
        l = Inches(0.4) + i * (metric_w + Inches(0.08))
        add_text(s, val, l, Inches(5.2), metric_w, Inches(0.65),
                 size=26, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
        add_text(s, label, l, Inches(5.85), metric_w, Inches(0.5),
                 size=11, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    add_text(s, "Track A  ·  Predictive Modelling Specialisation  ·  PlantIQ Team",
             Inches(0.5), Inches(6.95), Inches(12), Inches(0.4),
             size=12, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
def main():
    prs = new_prs()
    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_approach(prs)
    slide_04_functionality(prs)
    slide_05_key_features(prs)
    slide_06_architecture(prs)
    slide_07_ai_models(prs)
    slide_08_tech_stack(prs)
    slide_09_model_performance(prs)
    slide_10_realworld(prs)
    slide_11_closing(prs)

    out = "PlantIQ_Round2_Presentation.pptx"
    prs.save(out)
    print(f"✅  Saved  →  {out}  ({len(prs.slides)} slides)")

if __name__ == "__main__":
    main()
